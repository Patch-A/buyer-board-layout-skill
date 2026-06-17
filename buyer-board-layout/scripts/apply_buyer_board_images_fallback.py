from __future__ import annotations

import argparse
import json
import tempfile
from pathlib import Path
from typing import Any

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def remove_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


def normalize_image(path: Path, temp_dir: Path) -> Path:
    suffix = path.suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"}:
        return path

    if suffix == ".svg":
        import cairosvg

        out = temp_dir / f"{path.stem}.png"
        cairosvg.svg2png(url=str(path), write_to=str(out))
        return out

    if suffix == ".webp":
        out = temp_dir / f"{path.stem}.png"
        Image.open(path).save(out, "PNG")
        return out

    raise ValueError(f"Unsupported image format: {path}")


def find_shape(slide, left_pt: float, top_pt: float):
    target_left = Pt(left_pt)
    target_top = Pt(top_pt)
    tolerance = Pt(3)
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if abs(shape.left - target_left) <= tolerance and abs(shape.top - target_top) <= tolerance:
            return shape
    raise ValueError(f"Target picture shape not found near ({left_pt}, {top_pt})")


def fit_into_box(shape, left, top, width, height, fill: bool, align_left: bool = False) -> None:
    ratio_x = width / shape.width
    ratio_y = height / shape.height
    ratio = max(ratio_x, ratio_y) if fill else min(ratio_x, ratio_y)
    shape.width = int(shape.width * ratio)
    shape.height = int(shape.height * ratio)
    if align_left:
        shape.left = int(left)
    else:
        shape.left = int(left + ((width - shape.width) / 2))
    shape.top = int(top + ((height - shape.height) / 2))


def trim_uniform_border(image: Image.Image, tolerance: int = 12) -> Image.Image:
    rgba = image.convert("RGBA")
    bg = Image.new("RGBA", rgba.size, rgba.getpixel((0, 0)))
    diff = ImageChops.difference(rgba, bg)
    diff = ImageChops.add(diff, diff, 2.0, -tolerance)
    bbox = diff.getbbox()
    return rgba.crop(bbox) if bbox else rgba


def crop_to_aspect(image: Image.Image, target_width: int, target_height: int) -> Image.Image:
    source_ratio = image.width / image.height
    target_ratio = target_width / target_height
    if abs(source_ratio - target_ratio) < 0.01:
        return image

    if source_ratio > target_ratio:
        new_width = int(image.height * target_ratio)
        offset = (image.width - new_width) // 2
        return image.crop((offset, 0, offset + new_width, image.height))

    new_height = int(image.width / target_ratio)
    offset = (image.height - new_height) // 2
    return image.crop((0, offset, image.width, offset + new_height))


def prepare_site_fill_image(image_path: Path, temp_dir: Path, target_width: int, target_height: int) -> Path:
    normalized = normalize_image(image_path, temp_dir)
    image = Image.open(normalized).convert("RGB")
    cropped = crop_to_aspect(image, target_width, target_height)
    out = temp_dir / f"{image_path.stem}-site-fill.png"
    cropped.save(out, "PNG")
    return out


def prepare_logo_image(image_path: Path, temp_dir: Path) -> Path:
    normalized = normalize_image(image_path, temp_dir)
    image = Image.open(normalized)
    trimmed = trim_uniform_border(image)
    out = temp_dir / f"{image_path.stem}-logo.png"
    trimmed.save(out, "PNG")
    return out


def replace_picture(slide, target, image_path: Path, fill: bool, temp_dir: Path | None = None) -> None:
    left, top, width, height = target.left, target.top, target.width, target.height
    remove_shape(target)
    asset_to_use = image_path
    if fill:
        if temp_dir is None:
            raise ValueError("temp_dir is required when fill=True")
        asset_to_use = prepare_site_fill_image(image_path, temp_dir, width, height)
        slide.shapes.add_picture(str(asset_to_use), left, top, width=width, height=height)
        return

    new_shape = slide.shapes.add_picture(str(asset_to_use), left, top)
    fit_into_box(new_shape, left, top, width, height, False, align_left=True)


def clear_region(slide, region: dict[str, float]) -> None:
    left = Pt(region["left"])
    top = Pt(region["top"])
    right = Pt(region["right"])
    bottom = Pt(region["bottom"])
    to_delete = []
    for shape in slide.shapes:
        if shape.left >= left and shape.left <= right and shape.top >= top and shape.top <= bottom:
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PLACEHOLDER}:
                to_delete.append(shape)
    for shape in to_delete:
        remove_shape(shape)


def add_logo(slide, image_path: Path, logo_cfg: dict[str, float]) -> None:
    left = Pt(logo_cfg["left"])
    top = Pt(logo_cfg["top"])
    width = Pt(logo_cfg["width"])
    height = Pt(logo_cfg["height"])
    new_shape = slide.shapes.add_picture(str(image_path), left, top)
    fit_into_box(new_shape, left, top, width, height, False, align_left=True)


def main() -> int:
    parser = argparse.ArgumentParser(description="Fallback image application without PowerPoint COM.")
    parser.add_argument("--input-ppt", required=True)
    parser.add_argument("--buyers-json", required=True)
    parser.add_argument("--layout-config", required=True)
    parser.add_argument("--output-ppt", required=True)
    parser.add_argument("--preview-dir")
    args = parser.parse_args()

    input_ppt = Path(args.input_ppt)
    buyers = load_json(Path(args.buyers_json))
    layout = load_json(Path(args.layout_config))
    prs = Presentation(str(input_ppt))

    temp_dir = Path(tempfile.mkdtemp(prefix="buyer-board-images-"))
    start_index = int(layout["content"]["start_slide_index"]) - 1
    slots = {int(item["slide_offset"]): item for item in layout["images"]["slides"]}

    for idx, buyer in enumerate(buyers):
        slide = prs.slides[start_index + idx]
        slot = slots.get(idx)
        if not slot:
            continue

        if slot.get("site") and buyer.get("site_image_path"):
            site_cfg = slot["site"]
            site_asset = Path(buyer["site_image_path"])
            target = find_shape(slide, float(site_cfg["target_left"]), float(site_cfg["target_top"]))
            replace_picture(slide, target, site_asset, bool(site_cfg.get("fill", False)), temp_dir=temp_dir)

        if slot.get("logo") and buyer.get("logo_path"):
            logo_cfg = slot["logo"]
            logo_asset = prepare_logo_image(Path(buyer["logo_path"]), temp_dir)
            if logo_cfg["mode"] == "add":
                if logo_cfg.get("clear_region"):
                    clear_region(slide, logo_cfg["clear_region"])
                add_logo(slide, logo_asset, logo_cfg)
            else:
                target = find_shape(slide, float(logo_cfg["target_left"]), float(logo_cfg["target_top"]))
                replace_picture(slide, target, logo_asset, False, temp_dir=temp_dir)

    output_ppt = Path(args.output_ppt)
    output_ppt.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_ppt))
    if args.preview_dir:
        preview_dir = Path(args.preview_dir)
        preview_dir.mkdir(parents=True, exist_ok=True)
        (preview_dir / "README.txt").write_text(
            "Slide preview export was skipped because PowerPoint COM was unavailable. "
            "The PPT was generated with the Python fallback image pipeline.",
            encoding="utf-8",
        )
    print(output_ppt)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
